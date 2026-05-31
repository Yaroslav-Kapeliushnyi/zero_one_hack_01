"""
Build the full 10-slide presentation as a .pptx file.
Usage: python src/eval/make_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT  = Path(__file__).parent.parent.parent
PLOTS = ROOT / "plots"
OUT   = ROOT / "Zheng_et_al_slides.pptx"

# ── Colours ───────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)  # blue
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUBTEXT   = RGBColor(0x94, 0xA3, 0xB8)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_image(slide, path, left, top, width, height=None):
    if not Path(path).exists():
        return
    if height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width))


def section_title(slide, title, subtitle=None):
    """Top accent bar + slide title."""
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
    add_text(slide, title, 0.5, 0.18, 12.3, 0.7,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.88, 12.3, 0.45,
                 size=14, color=SUBTEXT)


def bullet_block(slide, items, left, top, width, height,
                 size=13, color=WHITE, spacing=0.38):
    y = top
    for item in items:
        bullet = "▸  " if not item.startswith("   ") else "     "
        add_text(slide, bullet + item.lstrip(), left, y, width, spacing,
                 size=size, color=color)
        y += spacing


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide1(prs):
    s = blank_slide(prs)
    bg(s)

    # Accent bar left
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)
    # Bottom bar
    add_rect(s, 0, 6.9, 13.33, 0.6, BG_CARD)

    add_text(s, "Zero One Hack 2026", 0.4, 0.4, 12, 0.6,
             size=14, color=SUBTEXT)
    add_text(s, "Learning Semiconductor\nProcess Logic", 0.4, 1.0, 11, 2.0,
             size=42, bold=True, color=WHITE)
    add_text(s, "with Neural Sequence Models", 0.4, 2.85, 11, 0.7,
             size=24, color=ACCENT)

    add_text(s, "Industrial AI Track  ·  Infineon Technologies", 0.4, 3.7, 9, 0.5,
             size=13, color=SUBTEXT)

    # Team block
    add_rect(s, 0.4, 4.3, 5.5, 2.1, BG_CARD)
    add_text(s, "Team Zheng et al", 0.65, 4.38, 5, 0.45,
             size=13, bold=True, color=ACCENT)
    for i, (name, role) in enumerate([
        ("Yehor Larcenko",        "ML lead · LSTM · evaluation"),
        ("Olga Rybak",            "GPT Transformer · infrastructure"),
        ("Yaroslav Kapeliushnyi", "Markov · metrics · TCN · GRPO"),
    ]):
        add_text(s, name, 0.65, 4.88 + i*0.45, 2.2, 0.4, size=11, bold=True, color=WHITE)
        add_text(s, role, 2.9,  4.88 + i*0.45, 3.2, 0.4, size=11, color=SUBTEXT)

    # Stack block
    add_rect(s, 6.5, 4.3, 6.4, 2.1, BG_CARD)
    add_text(s, "Cluster & Stack", 6.75, 4.38, 5.5, 0.45,
             size=13, bold=True, color=ACCENT)
    for i, line in enumerate([
        "Leonardo HPC  ·  NVIDIA A100-SXM-64GB",
        "Python 3.11  ·  PyTorch 2.5",
        "HuggingFace Transformers",
        "Custom eval pipeline (eval_metrics.py)",
    ]):
        add_text(s, "▸  " + line, 6.75, 4.85 + i*0.38, 5.9, 0.38,
                 size=10.5, color=WHITE)

    add_text(s, "Track: Industrial AI  ·  Learning & Benchmarking Process Logic",
             0.4, 6.95, 12, 0.45, size=11, color=SUBTEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem
# ─────────────────────────────────────────────────────────────────────────────
def slide2(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "The Problem", "What does the model need to learn?")

    # Sequence example
    add_rect(s, 0.4, 1.2, 12.5, 1.05, BG_CARD)
    add_text(s, "Example: IC chip manufacturing sequence (~130 steps)",
             0.6, 1.25, 12, 0.38, size=11, color=SUBTEXT)
    seq = "RECEIVE WAFER LOT  →  PRE CLEAN  →  THERMAL OXIDATION  →  SPIN COAT PHOTORESIST  →  EXPOSE LITHO  →  DEVELOP  →  OXIDE ETCH  →  …  →  SHIP LOT"
    add_text(s, seq, 0.6, 1.58, 12.1, 0.55, size=10.5, color=WHITE)

    # 3 task cards
    tasks = [
        ("01", "Next-Step Prediction",
         "Given 60% or 80% of a sequence\n→ predict the correct next step\n\nMetric: Top-1 / Top-3 / Top-5 Accuracy, MRR",
         ACCENT),
        ("02", "Sequence Completion",
         "Given 60% or 80% of a sequence\n→ generate ALL remaining steps\n\nMetric: Normalized Edit Distance, Token Accuracy",
         GREEN),
        ("03", "Anomaly Detection",
         "Given a complete sequence\n→ does it violate any of 10 process rules?\n\nMetric: F1, ROC-AUC, Rule Attribution",
         ORANGE),
    ]
    for i, (num, title, body, color) in enumerate(tasks):
        lft = 0.4 + i * 4.3
        add_rect(s, lft, 2.45, 4.1, 3.8, BG_CARD)
        add_rect(s, lft, 2.45, 4.1, 0.08, color)
        add_text(s, num, lft + 0.15, 2.52, 0.6, 0.55,
                 size=28, bold=True, color=color)
        add_text(s, title, lft + 0.15, 3.05, 3.7, 0.5,
                 size=14, bold=True, color=WHITE)
        add_text(s, body, lft + 0.15, 3.55, 3.7, 2.5,
                 size=11, color=SUBTEXT)

    add_text(s, "198 unique step names  ·  3 product families (MOSFET / IGBT / IC)  ·  10 formal grammar rules  ·  115–150 steps per sequence",
             0.4, 6.3, 12.5, 0.5, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Data
# ─────────────────────────────────────────────────────────────────────────────
def slide3(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Data & Setup", "From 3,000 to 33,000 training sequences")

    add_image(s, PLOTS / "05_data_volume.png", 0.4, 1.25, 7.8)

    add_rect(s, 8.7, 1.25, 4.3, 5.8, BG_CARD)
    add_text(s, "Key design decisions", 8.9, 1.35, 3.9, 0.45,
             size=13, bold=True, color=ACCENT)

    bullets = [
        "1K sequences/family provided",
        "Generated 9K more per family\nusing official generator script",
        "Total: 33,000 sequences",
        "",
        "Family conditioning token",
        "[MOSFET] / [IGBT] / [IC]\nprepended to every sequence",
        "→ one model for all families",
        "",
        "Vocabulary: 206 tokens",
        "198 process steps + 8 special",
        "",
        "Split: 90% train / 10% val",
        "Fixed seed=42 for reproducibility",
    ]
    y = 1.9
    for b in bullets:
        if b == "":
            y += 0.18
            continue
        color = WHITE if not b.startswith("→") else GREEN
        add_text(s, ("▸  " if not b.startswith(" ") and not b.startswith("→") else "   ") + b,
                 8.9, y, 3.8, 0.38, size=10.5, color=color)
        y += 0.38


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Baselines
# ─────────────────────────────────────────────────────────────────────────────
def slide4(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Baselines", "Three levels of baselines — showing where neural models add value")

    # Progression: Random → Unigram → GBM → Trigram → LSTM
    baselines = [
        ("Naive\nmost-freq","2.3%", "Always predict the globally most-frequent step",            RGBColor(0xEF,0x44,0x44)),
        ("Markov order-1\n(Unigram)","~47%", "Given last step → predict most common follower",   RGBColor(0xF9,0x73,0x16)),
        ("GBM\n(last 5 steps)","~59%","Gradient Boosted classifier, features = last 5 step IDs", RGBColor(0xF9,0x73,0x16)),
        ("Markov order-3\n(Trigram)","70.0%","Given last 3 steps → count frequencies in 33K seqs",RGBColor(0xEA,0xB3,0x08)),
        ("Our best\n(Hybrid)", "71.8%","Ensemble top pick + synonym-variant fill",              GREEN),
    ]

    for i, (label, acc, desc, color) in enumerate(baselines):
        left = 0.3 + i * 2.6
        add_rect(s, left, 1.25, 2.4, 4.8, BG_CARD)
        add_rect(s, left, 1.25, 2.4, 0.06, color)
        add_text(s, label,  left+0.12, 1.35, 2.18, 0.65,
                 size=12, bold=True, color=WHITE)
        add_text(s, acc,    left+0.12, 2.10, 2.18, 0.55,
                 size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, "Top-1", left+0.12, 2.65, 2.18, 0.30,
                 size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)
        add_text(s, desc,   left+0.12, 3.10, 2.18, 1.5,
                 size=10, color=SUBTEXT)

    # Arrow showing progression
    add_text(s, "─────────────────────────────────────────── Increasing complexity ───────────────────────────────────────────▶",
             0.3, 6.15, 12.7, 0.38, size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)

    add_text(s,
             "Key insight: Markov order-3 already reaches 70.0% by counting patterns — on identical "
             "examples it ties the ensemble. Neural capacity pays off on Task 2 completion (NED −38% "
             "vs trigram), not Task 1 Top-1.",
             0.3, 6.6, 12.7, 0.75, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Models trained
# ─────────────────────────────────────────────────────────────────────────────
def slide5(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Models Trained from Scratch", "All trained on Leonardo A100 — no pre-trained weights used for Tasks 1+2")

    add_image(s, PLOTS / "01_loss_curves.png", 0.4, 1.2, 8.5)

    models = [
        ("LSTM",  "2-layer, hidden=512", "3M",   "0.3293", ACCENT),
        ("GPT",   "8L / 8H / d=256",    "6.4M", "0.3287", RGBColor(0xA8,0x55,0xF7)),
        ("TCN",   "4-block dilated",     "~2M",  "0.3401", RGBColor(0x06,0xB6,0xD4)),
        ("Qwen\nGRPO", "1.5B + RL",     "1.5B", "—",      ORANGE),
    ]

    add_rect(s, 9.2, 1.2, 3.8, 5.8, BG_CARD)
    add_text(s, "Architecture", 9.4, 1.28, 3.4, 0.38, size=12, bold=True, color=ACCENT)
    add_text(s, "Params  Val Loss", 10.6, 1.28, 2.4, 0.38, size=10, color=SUBTEXT)

    for i, (name, desc, params, loss, color) in enumerate(models):
        y = 1.78 + i * 1.3
        add_rect(s, 9.25, y, 3.7, 1.15, RGBColor(0x0F,0x17,0x2A))
        add_rect(s, 9.25, y, 0.05, 1.15, color)
        add_text(s, name, 9.38, y + 0.08, 1.5, 0.5, size=13, bold=True, color=color)
        add_text(s, desc, 9.38, y + 0.55, 2.2, 0.38, size=10, color=SUBTEXT)
        add_text(s, params, 11.05, y + 0.08, 1.0, 0.38, size=12, bold=True, color=WHITE)
        add_text(s, loss,   11.05, y + 0.55, 1.0, 0.38, size=11, color=GREEN)

    add_text(s, "Key finding: LSTM and GPT converge to nearly identical val loss (~0.33).\nBottleneck is data structure, not model capacity.",
             0.4, 6.9, 12.5, 0.55, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Task 1 results
# ─────────────────────────────────────────────────────────────────────────────
def slide6(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Task 1 — Next-Step Prediction", "Top-1 / Top-3 / Top-5 Accuracy, MRR")

    add_image(s, PLOTS / "02_task1_models.png", 0.4, 1.2, 8.2)

    add_rect(s, 8.9, 1.2, 4.1, 5.8, BG_CARD)
    add_text(s, "Best result", 9.1, 1.28, 3.7, 0.4, size=13, bold=True, color=ACCENT)

    stats = [
        ("Top-1", "71.8%", GREEN),
        ("Top-3", "99.8%", GREEN),
        ("Top-5", "100%",  GREEN),
        ("MRR",   "0.857", GREEN),
    ]
    for i, (label, val, color) in enumerate(stats):
        y = 1.82 + i * 0.72
        add_rect(s, 9.1, y, 3.7, 0.6, RGBColor(0x0F,0x17,0x2A))
        add_text(s, label, 9.25, y + 0.1, 1.5, 0.4, size=12, color=SUBTEXT)
        add_text(s, val,   10.7, y + 0.08, 2.0, 0.45,
                 size=18, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_rect(s, 8.9, 4.75, 4.1, 2.0, RGBColor(0x0F,0x17,0x2A))
    add_text(s, "Hybrid method:", 9.1, 4.82, 3.7, 0.38,
             size=12, bold=True, color=WHITE)
    add_text(s,
             "RANK_1: LSTM-Attn + GPT + Markov\nensemble picks the top step\n"
             "RANK_2-5: synonym variants (lookup)\n\n"
             "Best rank-1 picker + full Top-5\ncoverage → beats every single model",
             9.1, 5.22, 3.7, 1.45, size=10.5, color=SUBTEXT)

    add_text(s, "Top-5 = 100%: the correct answer is ALWAYS in our top 5",
             0.4, 7.05, 12.5, 0.38, size=11.5, bold=True,
             color=GREEN, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Task 2 results
# ─────────────────────────────────────────────────────────────────────────────
def slide7(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Task 2 — Sequence Completion", "Normalized Edit Distance · Token Accuracy · Exact Match")

    add_image(s, PLOTS / "03_task2_ned.png", 0.4, 1.2, 8.8)

    add_rect(s, 9.5, 1.2, 3.5, 5.8, BG_CARD)
    add_text(s, "Best result", 9.7, 1.28, 3.1, 0.4, size=13, bold=True, color=ACCENT)

    stats2 = [
        ("NED ↓",       "0.2223", GREEN),
        ("Token Acc",   "42.0%",  GREEN),
        ("Block Acc",   "53.8%",  GREEN),
        ("Exact Match", "0.67%",  SUBTEXT),
    ]
    for i, (label, val, color) in enumerate(stats2):
        y = 1.82 + i * 0.72
        add_rect(s, 9.7, y, 3.1, 0.6, RGBColor(0x0F,0x17,0x2A))
        add_text(s, label, 9.85, y + 0.1, 1.5, 0.4, size=12, color=SUBTEXT)
        add_text(s, val, 11.0, y + 0.08, 1.6, 0.45,
                 size=18, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_rect(s, 9.5, 4.75, 3.5, 2.1, RGBColor(0x0F,0x17,0x2A))
    add_text(s, "Beam search (width=5):", 9.7, 4.82, 3.1, 0.38,
             size=12, bold=True, color=WHITE)
    add_text(s,
             "Keep 5 parallel paths at each step.\n"
             "Explore multiple possibilities,\n"
             "pick the best complete sequence.\n\n"
             "Greedy NED:  0.2245\n"
             "Beam-5 NED:  0.2223  ✓",
             9.7, 5.22, 3.1, 1.55, size=10.5, color=SUBTEXT)

    add_text(s, "NED = 0.22 means our generated suffix is ~78% correct token-by-token",
             0.4, 7.05, 12.5, 0.38, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Task 3 results
# ─────────────────────────────────────────────────────────────────────────────
def slide8(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Task 3 — Anomaly Detection", "Neuro-symbolic hybrid: rule engine + LSTM uncertainty score")

    add_image(s, PLOTS / "04_task3_anomaly.png", 0.4, 1.2, 12.5)

    add_text(s,
             "Binary detection (IS_VALID):  formal rule checker — deterministic, uses all 10 known grammar rules     "
             "Anomaly score (0–1):  LSTM NLL — higher uncertainty = more anomalous     "
             "Note: F1=1.0 is perfect by construction — our detector IS the official 10-rule checker. "
             "Rule Attribution 87.3% is the realistic measure.",
             0.4, 6.45, 12.5, 0.9, size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Task 3 Self-test
# ─────────────────────────────────────────────────────────────────────────────
def slide9_selftest(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Task 3 — Honest Self-Test", "Our own anomaly set (not organizer data) — 300 valid + 300 anomalous")

    add_image(s, PLOTS / "07_task3_selftest.png", 0.4, 1.2, 12.5)

    add_rect(s, 0.4, 6.35, 12.5, 1.0, BG_CARD)
    add_text(s, "Why two evaluations?", 0.65, 6.42, 4.0, 0.35,
             size=12, bold=True, color=ACCENT)
    add_text(s,
             "F1=1.0 is perfect by construction — our detector uses the official validate_sequence() "
             "10-rule checker, so it catches every rule violation. The informative number is rule "
             "attribution (87.3%): naming the exact violated rule is harder — some rules are ambiguous "
             "(e.g. BACKSIDE confused with PAD_OPEN). Official 987-seq score: graded by organizers.",
             0.65, 6.78, 12.0, 0.52, size=10.5, color=SUBTEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
def slide10_arch(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Architecture — Neuro-Symbolic Hybrid + Results", "How the system works · final scores")

    # Left: architecture flow (compact)
    boxes = [
        (0.3, 1.3, 2.8, 0.85, "INPUT",        "Partial sequence\n(60%/80%)",        BG_CARD, ACCENT),
        (3.4, 1.3, 2.8, 0.85, "FAMILY TOKEN", "[MOSFET]/[IGBT]/[IC]",               BG_CARD, ORANGE),
        (6.5, 1.3, 2.8, 0.85, "LSTM",         "3M params · 33K train seqs",         BG_CARD, ACCENT),
        (0.3, 2.7, 2.8, 0.85, "RULE ENGINE",  "10 rules → block invalid → −∞",      BG_CARD, GREEN),
        (3.4, 2.7, 2.8, 0.85, "RANKED STEPS", "LSTM probs over valid steps only",   BG_CARD, GREEN),
        (6.5, 2.7, 2.8, 0.85, "BEAM SEARCH",  "Width=5 · length norm",              BG_CARD, ORANGE),
    ]
    outputs_flow = [
        (0.5,  4.1, 2.4, 0.75, "Task 1", "Top-1 = 71.8%",  ACCENT),
        (3.6,  4.1, 2.4, 0.75, "Task 2", "NED = 0.2223",    GREEN),
        (6.7,  4.1, 2.4, 0.75, "Task 3", "F1 = 1.000*",     ORANGE),
    ]
    for l, t, w, h, title, body, bg_c, ac in boxes:
        add_rect(s, l, t, w, h, bg_c)
        add_rect(s, l, t, w, 0.05, ac)
        add_text(s, title, l+0.1, t+0.06, w-0.15, 0.3,  size=10, bold=True, color=WHITE)
        add_text(s, body,  l+0.1, t+0.38, w-0.15, 0.42, size=9,  color=SUBTEXT)
    for l, t, w, h, title, body, color in outputs_flow:
        add_rect(s, l, t, w, h, RGBColor(0x0F,0x17,0x2A))
        add_rect(s, l, t, w, 0.05, color)
        add_text(s, title, l+0.1, t+0.06, w-0.15, 0.3,  size=11, bold=True, color=color)
        add_text(s, body,  l+0.1, t+0.38, w-0.15, 0.32, size=12, bold=True, color=WHITE)

    # Right: summary table + demo
    add_rect(s, 9.8, 1.3, 3.2, 3.5, BG_CARD)
    add_text(s, "Demo: before vs after", 10.0, 1.38, 2.8, 0.38,
             size=12, bold=True, color=ACCENT)
    add_text(s, "Input: RECEIVE WAFER LOT\n→ LOT IDENTIFICATION\n→ INITIAL WAFER INSPECTION → ?",
             10.0, 1.82, 2.8, 0.65, size=9.5, color=SUBTEXT)
    demo = [
        ("Random",   "PASSIVATION ETCH",       RGBColor(0xEF,0x44,0x44), "❌"),
        ("Markov",   "MEASURE INIT. GEOMETRY",  ORANGE, "✓"),
        ("Ensemble", "MEASURE INIT. GEOMETRY",  GREEN,  "✓✓"),
    ]
    y = 2.56
    for model, pred, color, mark in demo:
        add_rect(s, 9.9, y, 3.1, 0.7, RGBColor(0x0F,0x17,0x2A))
        add_text(s, model,  9.95, y+0.04, 1.3, 0.3, size=9, color=SUBTEXT)
        add_text(s, pred,   9.95, y+0.34, 2.0, 0.3, size=9, bold=True, color=color)
        add_text(s, mark,  12.3,  y+0.15, 0.6, 0.4, size=16, color=color)
        y += 0.75

    add_rect(s, 9.8, 4.85, 3.2, 2.1, RGBColor(0x0F,0x17,0x2A))
    add_text(s, "What we learned", 10.0, 4.93, 2.8, 0.35, size=11, bold=True, color=ACCENT)
    for i, l in enumerate(["Data structure > model size",
                            "Beam search helps Task 2",
                            "Rules + neural = best combo",
                            "Markov ≈ LSTM on Top-1 (structured data!)"]):
        add_text(s, "▸ " + l, 10.0, 5.35 + i*0.35, 2.8, 0.32, size=9.5, color=SUBTEXT)

    add_text(s, "Team Zheng et al · Zero One Hack 2026 · Industrial AI Track",
             0.3, 7.1, 13.0, 0.35, size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Summary + Demo
# ─────────────────────────────────────────────────────────────────────────────
def slide11_summary(prs):
    s = blank_slide(prs)
    bg(s)
    section_title(s, "Results Summary & Demo", "Baseline → trained model comparison")

    add_image(s, PLOTS / "06_summary_table.png", 0.4, 1.2, 7.8)

    add_rect(s, 8.5, 1.2, 4.5, 5.8, BG_CARD)
    add_text(s, "Demo: before vs after", 8.7, 1.28, 4.1, 0.4,
             size=13, bold=True, color=ACCENT)

    add_text(s, 'Input: "RECEIVE WAFER LOT → LOT IDENTIFICATION →\nINITIAL WAFER INSPECTION → ?"',
             8.7, 1.78, 4.1, 0.7, size=10, color=SUBTEXT)

    demo_rows = [
        ("Random baseline", "PASSIVATION ETCH", RGBColor(0xEF,0x44,0x44), "❌"),
        ("Markov n-gram",   "MEASURE INITIAL\nGEOMETRY",   ORANGE, "✓"),
        ("LSTM Ensemble",   "MEASURE INITIAL\nGEOMETRY",   GREEN,  "✓✓"),
    ]
    y = 2.58
    for model, pred, color, mark in demo_rows:
        add_rect(s, 8.55, y, 4.35, 0.88, RGBColor(0x0F,0x17,0x2A))
        add_text(s, model, 8.7,  y+0.08, 2.2, 0.35, size=10, color=SUBTEXT)
        add_text(s, pred,  8.7,  y+0.42, 2.5, 0.42, size=10.5, bold=True, color=color)
        add_text(s, mark,  11.3, y+0.22, 0.9, 0.45, size=18, color=color, align=PP_ALIGN.CENTER)
        y += 0.95

    add_rect(s, 8.5, 5.55, 4.5, 1.4, RGBColor(0x0F,0x17,0x2A))
    add_text(s, "What we learned", 8.7, 5.62, 4.1, 0.38,
             size=12, bold=True, color=ACCENT)
    lessons = [
        "Data structure > model size (Markov ≈ LSTM on Top-1)",
        "Beam search helps completion but marginally",
        "Symbolic rules + neural = best of both worlds",
        "GRPO promising but needs format tuning",
    ]
    y = 6.05
    for l in lessons:
        add_text(s, "▸  " + l, 8.7, y, 4.1, 0.32, size=9.5, color=SUBTEXT)
        y += 0.32

    add_text(s, "Thank you  ·  Team Zheng et al  ·  Zero One Hack 2026",
             0.4, 7.08, 12.5, 0.38, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    print("Building slides...")
    for fn, label in [
        (slide1,          "01 — Title"),
        (slide2,          "02 — Problem"),
        (slide3,          "03 — Data"),
        (slide4,          "04 — Baselines"),
        (slide5,          "05 — Models"),
        (slide6,          "06 — Task 1"),
        (slide7,          "07 — Task 2"),
        (slide8,          "08 — Task 3"),
        (slide9_selftest, "09 — Task 3 Self-Test"),
        (slide10_arch,    "10 — Architecture + Summary"),
    ]:
        fn(prs)
        print(f"  ✓ Slide — {label}")

    prs.save(OUT)
    print(f"\nSaved → {OUT}")

if __name__ == "__main__":
    build()
